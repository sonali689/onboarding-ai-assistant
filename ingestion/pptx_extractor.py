import os
from pptx import Presentation
from config import PPTX_FOLDER


def scan_pptx_files(root_folder: str = PPTX_FOLDER) -> list[str]:
    """Recursively find all .pptx files across all subfolders."""
    pptx_files = []
    for root, dirs, files in os.walk(root_folder):
        for f in files:
            if f.endswith(".pptx") or f.endswith(".PPTX"):
                pptx_files.append(os.path.join(root, f))
    print(f"Found {len(pptx_files)} PPTX files under {root_folder}")
    return sorted(pptx_files)


def extract_slide_data(pptx_path: str) -> list[dict]:
    """
    Open a PPTX and return a list of dicts, one per slide.
    Preserves Japanese and English text exactly as-is (UTF-8).
    """
    slides_data = []
    filename = os.path.basename(pptx_path)
    subfolder = os.path.basename(os.path.dirname(pptx_path))

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠️  Could not open {filename}: {e}")
        return []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # ── Extract all text shapes ──────────────────────────────────────────
        text_parts = []
        slide_title = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # First non-empty text box becomes the title
                if not slide_title:
                    slide_title = text
                text_parts.append(text)

        # ── Speaker notes ────────────────────────────────────────────────────
        notes = ""
        try:
            if slide.has_notes_slide:
                raw = slide.notes_slide.notes_text_frame.text.strip()
                if raw:
                    notes = raw
        except Exception:
            pass

        slides_data.append({
            "slide_number": slide_num,
            "slide_title":  slide_title,
            "text":         "\n".join(text_parts),
            "notes":        notes,
            "filename":     filename,
            "subfolder":    subfolder,
        })

    return slides_data